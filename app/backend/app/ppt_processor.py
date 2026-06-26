import json
import os
import shutil
import uuid
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
from pptx import Presentation

from .db import get_conn, utc_now
from .formula_layout import recognize_formula_regions
from .formula_ocr import recognize_formula_image
from .formula_parser import (
    build_formula_dsl,
    extract_formula_candidates,
    normalize_broken_formula_text,
    to_json,
    try_evaluate_simple_numbers,
    validate_formula_candidate,
)
from .formula_structure import formula_diagnostics, repair_formula_latex

DATA_DIR = Path(os.environ.get("DATA_DIR", "/data"))
UPLOAD_DIR = DATA_DIR / "uploads"
ASSET_DIR = DATA_DIR / "assets"


def ensure_dirs() -> None:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def save_upload(fileobj, filename: str) -> tuple[str, str]:
    ensure_dirs()
    doc_id = str(uuid.uuid4())
    safe_name = filename.replace("/", "_").replace("\\", "_")
    dest = UPLOAD_DIR / f"{doc_id}_{safe_name}"
    with open(dest, "wb") as f:
        shutil.copyfileobj(fileobj, f)
    return doc_id, str(dest)


def process_document(doc_id: str, file_path: str, filename: str, ocr_engine: str | None = None) -> dict[str, Any]:
    ext = Path(filename).suffix.lower()
    now = utc_now()
    with get_conn() as conn:
        conn.execute(
            "INSERT INTO documents(id, filename, file_type, storage_path, page_count, status, created_at, updated_at) VALUES(?,?,?,?,?,?,?,?)",
            (doc_id, filename, ext.lstrip("."), file_path, 0, "PROCESSING", now, now),
        )

    if ext in [".pptx"]:
        count = _process_pptx(doc_id, file_path, ocr_engine)
    elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
        count = _process_image(doc_id, file_path, ocr_engine)
    elif ext in [".pdf"]:
        count = _process_pdf(doc_id, file_path, ocr_engine)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

    with get_conn() as conn:
        conn.execute(
            "UPDATE documents SET page_count=?, status='DONE', updated_at=? WHERE id=?",
            (count, utc_now(), doc_id),
        )
    return {"document_id": doc_id, "page_count": count, "status": "DONE"}


def _insert_page(conn, doc_id: str, page_no: int, title: str | None, text: str | None, image_path: str | None = None) -> None:
    conn.execute(
        "INSERT OR REPLACE INTO pages(document_id,page_no,title,extracted_text,image_path,ocr_status,created_at) VALUES(?,?,?,?,?,?,?)",
        (doc_id, page_no, title, text or "", image_path, "DONE", utc_now()),
    )


def _insert_asset(conn, doc_id: str, page_no: int, asset_type: str, asset_path: str | None, raw_text: str | None, metadata: dict[str, Any] | None = None) -> None:
    conn.execute(
        "INSERT INTO page_assets(document_id,page_no,asset_type,asset_path,raw_text,metadata_json,created_at) VALUES(?,?,?,?,?,?,?)",
        (doc_id, page_no, asset_type, asset_path, raw_text, json.dumps(metadata or {}, ensure_ascii=False), utc_now()),
    )


def _insert_formula(conn, doc_id: str, page_no: int, item: dict[str, Any], source_type_override: str | None = None) -> int:
    latex = item.get("latex") or item.get("raw_text") or ""
    latex = repair_formula_latex(latex)
    latex = normalize_broken_formula_text(latex)
    item["latex"] = latex
    dsl = build_formula_dsl(latex)
    diagnostics = formula_diagnostics(latex)
    if diagnostics:
        dsl["structure_diagnostics"] = diagnostics
    variables = dsl.get("variables", [])
    status, msg = validate_formula_candidate(latex)
    if diagnostics.get("warnings"):
        msg = f"{msg}; structure_warnings={','.join(diagnostics['warnings'])}"
    numeric_hint = try_evaluate_simple_numbers(latex)
    if numeric_hint:
        dsl["numeric_hint"] = numeric_hint
    final_status = item.get("status") or ("CANDIDATE" if status == "PASS" and not diagnostics.get("warnings") else "NEEDS_REVIEW")
    cur = conn.execute(
        """
        INSERT INTO formula_blocks(
            document_id,page_no,formula_seq,formula_title,raw_text,latex,normalized_latex,
            formula_dsl_json,variables_json,confidence,source_type,status,bbox_json,
            validation_message,created_at,updated_at
        ) VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)
        """,
        (
            doc_id,
            page_no,
            item.get("formula_seq", 1),
            item.get("formula_title"),
            item.get("raw_text"),
            latex,
            item.get("normalized_latex") or latex,
            to_json(dsl),
            to_json(variables),
            float(item.get("confidence", 0.0)),
            source_type_override or item.get("source_type", "unknown"),
            final_status,
            json.dumps(item.get("bbox"), ensure_ascii=False) if item.get("bbox") else None,
            msg,
            utc_now(),
            utc_now(),
        ),
    )
    formula_id = int(cur.lastrowid)
    conn.execute(
        "INSERT INTO formula_validation_results(formula_id,rule_name,result,message,created_at) VALUES(?,?,?,?,?)",
        (formula_id, "basic_syntax", status, msg, utc_now()),
    )
    conn.execute(
        "INSERT INTO formula_validation_results(formula_id,rule_name,result,message,created_at) VALUES(?,?,?,?,?)",
        (formula_id, "structure_diagnostics", "PASS" if not diagnostics.get("warnings") else "WARN", to_json(diagnostics), utc_now()),
    )
    return formula_id


def _insert_structured_image_formulas(conn, doc_id: str, page_no: int, image_path: str, ocr_engine: str | None, start_seq: int, title: str) -> int:
    """Run formula-region OCR and insert all structured candidates."""
    inserted = 0
    try:
        candidates = recognize_formula_regions(image_path, ocr_engine)
    except Exception:
        candidates = []

    if not candidates:
        rec = recognize_formula_image(image_path, ocr_engine)
        if rec.get("latex") or rec.get("raw_text"):
            candidates = [
                {
                    "formula_seq": start_seq,
                    "formula_title": title,
                    "raw_text": rec.get("raw_text"),
                    "latex": rec.get("latex"),
                    "confidence": rec.get("confidence", 0.0),
                    "source_type": f"image_formula_{rec.get('engine')}",
                    "status": rec.get("status", "NEEDS_REVIEW"),
                }
            ]

    for offset, cand in enumerate(candidates):
        cand["formula_seq"] = start_seq + offset
        cand.setdefault("formula_title", title)
        _insert_formula(conn, doc_id, page_no, cand)
        inserted += 1
    return inserted


def _process_pptx(doc_id: str, file_path: str, ocr_engine: str | None = None) -> int:
    prs = Presentation(file_path)
    with get_conn() as conn:
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            tables: list[list[list[str]]] = []
            title = None
            formula_seq = 1
            slide_asset_dir = ASSET_DIR / doc_id / f"page_{idx}"
            slide_asset_dir.mkdir(parents=True, exist_ok=True)

            for shape_no, shape in enumerate(slide.shapes, start=1):
                if hasattr(shape, "text") and shape.text:
                    if title is None and shape.text.strip():
                        title = shape.text.strip().splitlines()[0][:120]
                    texts.append(shape.text)

                if getattr(shape, "has_table", False):
                    table_rows = []
                    for row in shape.table.rows:
                        table_rows.append([cell.text for cell in row.cells])
                    tables.append(table_rows)
                    _insert_asset(conn, doc_id, idx, "table", None, json.dumps(table_rows, ensure_ascii=False), {"shape_no": shape_no})

                if getattr(shape, "shape_type", None) and hasattr(shape, "image"):
                    try:
                        ext = shape.image.ext or "png"
                        img_path = slide_asset_dir / f"image_{shape_no}.{ext}"
                        with open(img_path, "wb") as f:
                            f.write(shape.image.blob)
                        inserted = _insert_structured_image_formulas(
                            conn,
                            doc_id,
                            idx,
                            str(img_path),
                            ocr_engine,
                            formula_seq,
                            "PPT 이미지 수식 구조 OCR 후보",
                        )
                        _insert_asset(conn, doc_id, idx, "image", str(img_path), f"structured_formula_candidates={inserted}", {"shape_no": shape_no, "engine": ocr_engine})
                        formula_seq += inserted
                    except Exception as exc:
                        _insert_asset(conn, doc_id, idx, "image_error", None, str(exc), {"shape_no": shape_no})

            all_text = "\n".join(texts)
            _insert_page(conn, doc_id, idx, title, all_text)
            for c in extract_formula_candidates(all_text):
                c["formula_seq"] = formula_seq
                _insert_formula(conn, doc_id, idx, c)
                formula_seq += 1
    return len(prs.slides)


def _process_image(doc_id: str, file_path: str, ocr_engine: str | None = None) -> int:
    page_no = 1
    asset_dir = ASSET_DIR / doc_id / "page_1"
    asset_dir.mkdir(parents=True, exist_ok=True)
    asset_path = asset_dir / Path(file_path).name
    shutil.copyfile(file_path, asset_path)
    with get_conn() as conn:
        _insert_page(conn, doc_id, page_no, "이미지 수식 페이지", "", str(asset_path))
        inserted = _insert_structured_image_formulas(
            conn,
            doc_id,
            page_no,
            str(asset_path),
            ocr_engine,
            1,
            "이미지 수식 구조 OCR 후보",
        )
        _insert_asset(conn, doc_id, page_no, "image", str(asset_path), f"structured_formula_candidates={inserted}", {"engine": ocr_engine})
    return 1


def _process_pdf(doc_id: str, file_path: str, ocr_engine: str | None = None) -> int:
    pdf = fitz.open(file_path)
    with get_conn() as conn:
        for page_idx, page in enumerate(pdf, start=1):
            text = page.get_text("text") or ""
            pix = page.get_pixmap(matrix=fitz.Matrix(1.8, 1.8))
            page_dir = ASSET_DIR / doc_id / f"page_{page_idx}"
            page_dir.mkdir(parents=True, exist_ok=True)
            img_path = page_dir / "page.png"
            pix.save(str(img_path))
            _insert_page(conn, doc_id, page_idx, f"PDF Page {page_idx}", text, str(img_path))
            formula_seq = 1
            for c in extract_formula_candidates(text):
                c["formula_seq"] = formula_seq
                _insert_formula(conn, doc_id, page_idx, c)
                formula_seq += 1

            # Always run structured region OCR for rendered PDF pages.
            # Text extraction alone cannot recover 2D formulas such as fractions,
            # summation bounds, integral bounds, and multi-line actuarial formulas.
            inserted = _insert_structured_image_formulas(
                conn,
                doc_id,
                page_idx,
                str(img_path),
                ocr_engine,
                formula_seq,
                "PDF 페이지 수식 영역 구조 OCR 후보",
            )
            _insert_asset(conn, doc_id, page_idx, "page_image", str(img_path), f"structured_formula_candidates={inserted}", {"engine": ocr_engine})
    return len(pdf)
