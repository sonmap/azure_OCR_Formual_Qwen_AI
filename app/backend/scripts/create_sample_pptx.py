from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

out = Path(__file__).resolve().parents[2] / "samples" / "sample_actuarial_formula.pptx"
out.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "보험계리식 샘플 - 순보험료"
box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.2), Inches(4.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "예상 총 지급 보험금 (Expected Total Claim): 1년 동안 사망할 것으로 예상되는 인원수를 구한 뒤 보장 금액을 곱합니다."
p.font.size = Pt(18)
for text in [
    r"$$\text{예상 사망자 수} = 100,000\text{명} \times 0.002 = 200\text{명}$$",
    r"$$\text{예상 총 지급액} = 200\text{명} \times 100,000,000\text{원} = 20,000,000,000\text{원}$$",
    "1인당 연간 순보험료 (Net Premium per Person):",
    r"$$\text{1인당 순보험료} = \frac{\text{총 지급액}}{\text{총 가입자 수}} = \frac{20,000,000,000\text{원}}{100,000\text{명}} = 200,000\text{원}$$",
]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(16)

slide2 = prs.slides.add_slide(prs.slide_layouts[5])
slide2.shapes.title.text = "위험률표 샘플"
table = slide2.shapes.add_table(4, 4, Inches(0.8), Inches(1.2), Inches(8.0), Inches(2.0)).table
headers = ["product_id", "coverage_id", "age", "mortality_rate"]
rows = [headers, ["P001", "DEATH_BASIC", "40", "0.002"], ["P001", "DEATH_BASIC", "41", "0.0021"], ["P001", "DEATH_BASIC", "42", "0.0023"]]
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        table.cell(r, c).text = val

prs.save(out)
print(out)
